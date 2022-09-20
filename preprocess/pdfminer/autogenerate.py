from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches
import argparse
import sys
from lxml import etree


# outfilename = "output.pptx"
# inputname='example.md'
# pagesnum=3



class autogenerate():

    def __init__(self,outfilename,inputname,pagesnum):
        self.outfilename = outfilename
        self.inputname=inputname
        self.pagesnum=pagesnum

    def inputFileRead(self,filename):
        try:
            f = open(filename, "r",encoding="UTF-8")
            fcontent = f.readline()
            return fcontent, f
        except IOError:
            print("Error! File inaccessible")
            sys.exit(1)


    def allmain(self):
        if self.inputname is not None:
            prs = Presentation()
            title_slide_layout = prs.slide_layouts[0]
            content_slide_layout = prs.slide_layouts[1]
            title_slide = prs.slides.add_slide(title_slide_layout)
            content_slide = []
            if self.pagesnum is None:
                print("pages flag should be number.")
                sys.exit(1)
            for _ in range(int(self.pagesnum) - 1):
                content_slide.append(prs.slides.add_slide(content_slide_layout))
            title = title_slide.shapes.title
            subtitle = title_slide.placeholders[1]
            draft, f = self.inputFileRead(str(self.inputname))
            
            self.createPPTX(title, subtitle, draft, f, prs, content_slide)
        else:
            print("error")


    def createPPTX(self,title, subtitle, draft, f, prs, content_slide):
        signCount = 0
        callCount = 0
        while draft:
            fileContent = str(draft.strip())
            if "#" in fileContent and not fileContent.lstrip().startswith('##'):
                try:
                    type(PCDraft)
                except NameError:
                    title.text = str(fileContent.replace("#", ""))
                    prs.save(self.outfilename)
                else:
                    try:
                        if(PCDraft != 0):
                            for shape in content_slide[PCDraft].shapes:
                                if not shape.has_text_frame:
                                    continue
                                text_frame = shape.text_frame
                                p = text_frame.paragraphs[0]
                                p.text = fileContent.replace("#", "")
                        content_title = content_slide[PCDraft].shapes.title
                        content_title.text = str(fileContent.replace("#", ""))
                    except(IndexError):
                        print(
                            "Pages is not equal to section in draft file.",
                            "exiting..."
                        )
                        sys.exit(1)
                    except(UnboundLocalError):
                        print(
                            "Slide did not have content pages but the draft ",
                            "file specify one. Existing."
                        )
                        sys.exit(1)
                    prs.save(self.outfilename)
            elif "##" in fileContent:
                subtitle.text = str(fileContent.replace("##", ""))
                prs.save(self.outfilename)
            elif "\\newpage" in fileContent:
                callCount += 1
                if self.pagesnum <= 1:
                    pass
                else:
                    if callCount > 1:
                        PCDraft = signCount + 1
                        signCount += 1
                    else:
                        PCDraft = 0

    #        elif "content_title>" in fileContent:
    #            try:
    #                content_title = content_slide[PCDraft].shapes.title
    #                content_title.text = str(fileContent.replace("%%%", ""))
    #            except(IndexError):
    #                print(
    #                    "Pages is not equal to section in draft file.",
    #                    "exiting..."
    #                )
    #                sys.exit(1)
    #            except(UnboundLocalError):
    #                print(
    #                    "Slide did not have content pages but the draft ",
    #                    "file specify one. Existing."
    #                )
    #                sys.exit(1)
    #            prs.save(outfilename)

            elif "p>" in fileContent:
                try:
    #                for shape in content_slide[PCDraft].shapes:
    #                    if not shape.has_text_frame:
    #                        continue
                        _, body, *_ = content_slide[PCDraft].shapes.placeholders
                        p = body.text_frame.add_paragraph()
                        p.text = fileContent.replace("p>", "")
                        p._pPr.insert(0, etree.Element("{http://schemas.openxmlformats.org/drawingml/2006/main}buNone"))
                except(IndexError):
                    print("Pages is not enough. exiting...")
                    sys.exit(1)
                prs.save(self.outfilename)
            elif "bgcl>" in fileContent:
                background = content_slide[PCDraft].background
                try:
                    R, G, B = str(fileContent.replace("bgcl>", "")).split(",")
                except(ValueError):
                    print(
                        "background value should be RGB value seperate ",
                        "by comma. Example: bgcl>255,255,255"
                    )
                    sys.exit(1)
                fill = background.fill
                fill.solid()
                try:
                    fill.fore_color.rgb = RGBColor(int(R), int(G), int(B))
                except(ValueError):
                    print(
                        "background value should be RGB value seperate by ",
                        "comma. Example: bgcl>255,255,255"
                    )
                    sys.exit(1)
                prs.save(self.outfilename)
            elif "-" in fileContent:
                try:
    #                for shape in content_slide[PCDraft].shapes:
    #                    if not shape.has_text_frame:
    #                        continue
    #                    text_frame = shape.text_frame
                        _, body, *_ = content_slide[PCDraft].shapes.placeholders
                        p = body.text_frame.add_paragraph()
                        p.level = 0
                        p.text = fileContent.replace("-","")
                except(IndexError):
                    print("Page is not enough, exiting...")
                    sys.exit(1)
                prs.save(self.outfilename)
                    
            elif "img>" in fileContent:
                con_shapes = content_slide[PCDraft].shapes
                try:
                    imgpath, x1, y1, x2, y2 = str(
                        fileContent.replace("img>",
                                            "")
                    ).split(",")
                except(ValueError):
                    print(
                        "image value should be image path, Position of ",
                        "image horizontal, Position of image verticle, ",
                        "size of image horizontal, size of image verticle in Inches."
                    )
                    sys.exit(1)
                try:
                    con_shapes.add_picture(
                        str(imgpath), Inches(int(x1)),
                        Inches(int(y1)), Inches(int(x2)),
                        Inches(int(y2))
                    )
                except(ValueError):
                    print(
                        "image value should be image path, Position of ",
                        "image horizontal, Position of image verticle, ",
                        "size of image horizontal, size of image verticle in Inches."
                    )
                    sys.exit(1)
                except(FileNotFoundError):
                    print("image file did not exist. Please try again.")
                    sys.exit(1)
                prs.save(self.outfilename)
            elif "\n" or "\r\n" in fileContent:
                pass
            else:
                print(
                    "draft file is corrupted. Check the format.txt file ",
                    "and try again."
                )
            draft = f.readline()



